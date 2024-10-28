import re
import asyncio
import os
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO
from app.expand_content import expand_content_batch
from app.utils import clean_text

BATCH_SIZE = 10

def process_slide(slide):
    slide_content = {
        'title': None,
        'subtitles': [],
        'bullet_points': [],
        'tables': [],
        'notes': None
    }

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and not text.isdigit():
                    if para.level == 0:
                        slide_content['title'] = slide_content['title'] or text
                    elif para.level == 1:
                        slide_content['bullet_points'].append(text)
                    else:
                        slide_content['subtitles'].append(text)

        elif shape.shape_type == 17:  # Table
            slide_content['tables'].append([[cell.text.strip() for cell in row.cells] for row in shape.table.rows])

    if slide.has_notes_slide:
        slide_content['notes'] = slide.notes_slide.notes_text_frame.text.strip()

    return slide_content

def process_slides_sync(prs):
    with ThreadPoolExecutor(max_workers=os.cpu_count() * 2) as executor:
        slides_content = list(executor.map(process_slide, prs.slides))
    return slides_content

async def extract_and_expand_content(pptx_data: BytesIO):
    prs = Presentation(pptx_data)
    slides_content = await asyncio.to_thread(process_slides_sync, prs)

    all_bullet_points = [point for slide in slides_content if slide['bullet_points'] for point in slide['bullet_points']]
    batches = [all_bullet_points[i:i + BATCH_SIZE] for i in range(0, len(all_bullet_points), BATCH_SIZE)]

    expanded_points = []
    for batch in batches:
        expanded_points.extend(await expand_content_batch(batch))

    point_index = 0
    for slide in slides_content:
        if slide['bullet_points']:
            slide['expanded_bullet_points'] = expanded_points[point_index:point_index + len(slide['bullet_points'])]
            point_index += len(slide['bullet_points'])

    return slides_content

def format_notes_to_json(slides_content):
    formatted_slides = []
    next_id = 1  # Start IDs from 1
    
    for i, slide in enumerate(slides_content):
        # Skip the slide if it only has a title but no other content
        if not (slide['bullet_points'] or slide['subtitles'] or slide['tables'] or slide['notes']):
            continue  # Skip the slide if it's empty

        # Begin formatting the slide content
        slide_html = f"<h3>{slide['title']}</h3>"
        
        if slide['subtitles']:
            slide_html += '<p><strong>Subtitles:</strong></p><ul>' + ''.join(f'<li>{subtitle}</li>' for subtitle in slide['subtitles']) + '</ul>'
        
        if slide['bullet_points']:
            slide_html += '<p><strong>Key Concepts:</strong></p>'
            ol_open = False  # Track ordered list status
            ul_open = False  # Track unordered list status

            # Iterate over the bullet points and their expanded content
            for point, expanded in zip(slide['bullet_points'], slide['expanded_bullet_points']):
                point = clean_text(re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', point))
                expanded = clean_text(expanded if isinstance(expanded, str) else '')
                expanded = clean_text(re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', expanded))

                # Check if the expanded content starts with bold, treat it as a sub-bullet
                if re.search(r'^<b>', expanded):
                    # This is a sub-bullet for the current bullet point
                    if ul_open:
                        slide_html += f'<li>{point}<ul><li>{expanded}</li></ul></li>'
                    else:
                        slide_html += f'<ul><li>{point}<ul><li>{expanded}</li></ul></li></ul>'
                        ul_open = True
                else:
                    # Check if this is a numbered list item (e.g., "1.")
                    if re.match(r'^\d+\.', point):
                        if not ol_open:
                            slide_html += '<ol>'
                            ol_open = True
                        if ul_open:
                            slide_html += '</ul>'  # Close unordered list if one is open
                            ul_open = False
                        slide_html += f'<li>{point}<ul><li>{expanded}</li></ul></li>'
                    # Handle sub-bullets (starts with `-`, `•`, or bold sub-points)
                    elif point.startswith('-') or point.startswith('•') or re.match(r'^<b>', point):
                        if not ul_open:
                            slide_html += '<ul>'  # Open unordered list if one is not open
                            ul_open = True
                        if ol_open:
                            slide_html += '</ol>'  # Close ordered list if one is open
                            ol_open = False
                        # Add sub-bullet with margin for indentation
                        slide_html += f'<li>{point}</li><li style="margin-left:20px;">{expanded}</li>'
                    else:
                        # Handle normal bullets without sub-bullet characters
                        if ol_open:
                            slide_html += '</ol>'  # Close ordered list
                            ol_open = False
                        if ul_open:
                            slide_html += '</ul>'  # Close unordered list
                            ul_open = False
                        slide_html += f'<ul><li>{point}</li><li style="margin-left:20px;">{expanded}</li></ul>'

            # Ensure any open lists are closed
            if ol_open:
                slide_html += '</ol>'
            if ul_open:
                slide_html += '</ul>'
        
        if slide['tables']:
            slide_html += '<p><strong>Tables:</strong></p><ul>' + ''.join(f'<li>Table Data:<ul>' + ''.join(f'<li>{", ".join(row)}</li>' for row in table) + '</ul></li>' for table in slide['tables']) + '</ul>'
        
        if slide['notes']:
            slide_html += '<p><strong>Notes:</strong></p>'
            notes = clean_text(slide['notes'])

            # Split the notes into sentences or sections for better formatting
            note_lines = re.split(r'(?<=[.!?]) +|\n', notes)  # Split based on sentence boundaries or newlines

            # Apply similar formatting as done for bullet points
            slide_html += '<ul>'
            for line in note_lines:
                if re.search(r'^-|\*|•', line):  # If the note starts with a bullet-like character
                    slide_html += f'<ul><li>{line}</li></ul>'
                else:
                    slide_html += f'<li>{line}</li>'
            slide_html += '</ul>'
        
        # Append the formatted slide to the list as JSON
        formatted_slides.append({
            'id': next_id,
            'title': slide['title'],
            'content': slide_html
        })
        next_id += 1  # Increment the ID for the next slide

    return formatted_slides